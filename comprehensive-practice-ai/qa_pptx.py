"""用 python-pptx 检查越界，并为每页生成几何预览。"""
import json
import sys
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

src = Path(sys.argv[1])
out = Path(sys.argv[2])
out.mkdir(parents=True, exist_ok=True)
prs = Presentation(src)
sw, sh = prs.slide_width, prs.slide_height
scale = 96 / 914400
font = ImageFont.load_default()
issues = []

for i, slide in enumerate(prs.slides, 1):
    canvas = Image.new('RGB', (1280, 720), '#F7FAFF')
    draw = ImageDraw.Draw(canvas)
    for j, shape in enumerate(slide.shapes, 1):
        x, y, w, h = shape.left, shape.top, shape.width, shape.height
        if x < 0 or y < 0 or x + w > sw or y + h > sh:
            issues.append({'page': i, 'shape': j, 'name': shape.name, 'x': x, 'y': y, 'w': w, 'h': h})
        box = [round(x * scale), round(y * scale), round((x + w) * scale), round((y + h) * scale)]
        draw.rectangle(box, outline='#8A96A5', width=1)
        if getattr(shape, 'has_text_frame', False) and shape.text.strip():
            text = shape.text.replace('\n', ' ')[:90]
            draw.text((box[0] + 3, box[1] + 2), text, fill='#27364A', font=font)
    canvas.save(out / f'{i:03d}.png')

(out / 'report.json').write_text(json.dumps({'slides': len(prs.slides), 'issues': issues}, ensure_ascii=False, indent=2), encoding='utf-8')
print(f'{src.name}：{len(prs.slides)} 页，越界形状 {len(issues)} 个')
sys.exit(2 if issues else 0)
